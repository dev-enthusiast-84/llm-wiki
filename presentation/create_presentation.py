"""
PowerPoint/PPTX Generator for LLM Wiki Presentation
This script generates a complete PowerPoint presentation with all 5 slides.
Install: pip install python-pptx
Run: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def add_gradient_fill(shape, color1, color2):
    """Add gradient fill to a shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = RGBColor(*color1)
    fill.gradient_stops[1].color.rgb = RGBColor(*color2)

def set_text_color(text_frame, color):
    """Set text color for all paragraphs in a text frame"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(*color)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color schemes for each slide
colors = {
    1: {'bg': (102, 126, 234), 'accent': (118, 75, 162), 'text': (255, 255, 255)},  # Purple gradient
    2: {'bg': (240, 147, 251), 'accent': (245, 87, 108), 'text': (255, 255, 255)},   # Pink gradient
    3: {'bg': (79, 172, 254), 'accent': (0, 242, 254), 'text': (51, 51, 51)},        # Blue gradient
    4: {'bg': (250, 112, 154), 'accent': (254, 225, 64), 'text': (51, 51, 51)},      # Orange gradient
    5: {'bg': (48, 207, 208), 'accent': (51, 8, 103), 'text': (255, 255, 255)},      # Cyan gradient
}

# ============= SLIDE 1: TITLE SLIDE =============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.gradient()
fill.gradient_angle = 45.0
fill.gradient_stops[0].color.rgb = RGBColor(102, 126, 234)
fill.gradient_stops[1].color.rgb = RGBColor(118, 75, 162)

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.text = "🧠 From Concept to Implementation"
title_para.font.size = Pt(60)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = "Building an LLM Wiki: A Smart Knowledge Architecture"
subtitle_para.font.size = Pt(32)
subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
subtitle_para.alignment = PP_ALIGN.CENTER

# Tagline
tagline_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
tagline_frame = tagline_box.text_frame
tagline_para = tagline_frame.paragraphs[0]
tagline_para.text = "Turning Andrej Karpathy's Vision into Reality"
tagline_para.font.size = Pt(18)
tagline_para.font.color.rgb = RGBColor(255, 255, 255)
tagline_para.alignment = PP_ALIGN.CENTER

# Speaker notes
notes_slide = slide1.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = """Welcome everyone! Today we're exploring an innovative approach to managing knowledge using LLMs. 

We'll walk through Karpathy's concept of an LLM-powered wiki and show you how it's been brought to life in a real implementation. This 10-minute journey will transform how you think about personal knowledge bases.

Key points to emphasize:
- Knowledge management is broken for researchers and engineers
- Karpathy proposed an AI-powered solution
- We've built a working implementation with 3 Copilot Skills
- This demo shows what's possible with LLM-assisted knowledge work"""

# ============= SLIDE 2: THE PROBLEM & CONCEPT =============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.gradient()
fill.gradient_angle = 45.0
fill.gradient_stops[0].color.rgb = RGBColor(240, 147, 251)
fill.gradient_stops[1].color.rgb = RGBColor(245, 87, 108)

# Title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame2 = title_box2.text_frame
title_para2 = title_frame2.paragraphs[0]
title_para2.text = "⚡ The Problem We're Solving"
title_para2.font.size = Pt(48)
title_para2.font.bold = True
title_para2.font.color.rgb = RGBColor(255, 255, 255)

# Problem icons - row 1
problems = [
    ("📚", "Knowledge Overload"),
    ("🔗", "Scattered Sources"),
    ("😵", "No Structure"),
]

x_pos = 1
for icon, label in problems:
    icon_box = slide2.shapes.add_textbox(Inches(x_pos), Inches(1.8), Inches(2.5), Inches(0.6))
    icon_frame = icon_box.text_frame
    icon_para = icon_frame.paragraphs[0]
    icon_para.text = icon
    icon_para.font.size = Pt(50)
    icon_para.alignment = PP_ALIGN.CENTER
    
    label_box = slide2.shapes.add_textbox(Inches(x_pos), Inches(2.6), Inches(2.5), Inches(0.5))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label
    label_para.font.size = Pt(16)
    label_para.font.bold = True
    label_para.font.color.rgb = RGBColor(255, 255, 255)
    label_para.alignment = PP_ALIGN.CENTER
    
    x_pos += 2.8

# Solution box
solution_shape = slide2.shapes.add_shape(1, Inches(0.75), Inches(3.8), Inches(8.5), Inches(2.8))
solution_shape.fill.solid()
solution_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
solution_shape.line.color.rgb = RGBColor(245, 87, 108)
solution_shape.line.width = Pt(3)

solution_box = slide2.shapes.add_textbox(Inches(1), Inches(3.95), Inches(8), Inches(2.5))
solution_frame = solution_box.text_frame
solution_frame.word_wrap = True

# Heading
heading = solution_frame.paragraphs[0]
heading.text = "✨ Karpathy's Insight"
heading.font.size = Pt(28)
heading.font.bold = True
heading.font.color.rgb = RGBColor(245, 87, 108)
heading.level = 0

# Content
content = solution_frame.add_paragraph()
content.text = "Use AI to transform raw research into interconnected knowledge"
content.font.size = Pt(20)
content.font.bold = True
content.font.color.rgb = RGBColor(51, 51, 51)
content.level = 0
content.space_before = Pt(10)

detail = solution_frame.add_paragraph()
detail.text = "Papers become structured entities, concepts link together, and your wiki grows intelligently."
detail.font.size = Pt(18)
detail.font.color.rgb = RGBColor(80, 80, 80)
detail.level = 0
detail.space_before = Pt(8)

# Speaker notes for slide 2
notes_slide2 = slide2.notes_slide
text_frame2 = notes_slide2.notes_text_frame
text_frame2.text = """Every researcher, engineer, and student faces information overload. We read papers, articles, and docs, but they stay isolated. Karpathy proposed using LLMs not just to answer questions, but to systematically build and maintain knowledge bases.

The key insight: let AI extract concepts, create pages, and link ideas. This compounds knowledge growth over time.

Talk about:
- How traditional note-taking is scattered
- The pain of manually linking concepts
- Why AI-assisted knowledge graphs are powerful
- The long-term compounding effect of interconnected knowledge"""

# ============= SLIDE 3: THREE-STEP ARCHITECTURE =============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.gradient()
fill.gradient_angle = 45.0
fill.gradient_stops[0].color.rgb = RGBColor(79, 172, 254)
fill.gradient_stops[1].color.rgb = RGBColor(0, 242, 254)

# Title
title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
title_frame3 = title_box3.text_frame
title_para3 = title_frame3.paragraphs[0]
title_para3.text = "🏗️ The Three-Step Architecture"
title_para3.font.size = Pt(48)
title_para3.font.bold = True
title_para3.font.color.rgb = RGBColor(51, 51, 51)

# Four steps in a grid
steps = [
    ("1️⃣", "Compile", "Extract concepts from\npapers → Create pages"),
    ("2️⃣", "Sync", "Add new sources →\nUpdate & link pages"),
    ("3️⃣", "Audit", "Quality checks →\nFix orphans & contradictions"),
    ("🔄", "Repeat", "Continuously grow &\nimprove your wiki"),
]

y_start = 1.4
x_positions = [0.75, 5.25]
row = 0

for idx, (icon, title, desc) in enumerate(steps):
    x_pos = x_positions[idx % 2]
    y_pos = y_start + (row * 2.5)
    if idx == 2:
        row = 1
    
    # Box
    step_shape = slide3.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(4), Inches(2))
    step_shape.fill.solid()
    step_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    step_shape.line.color.rgb = RGBColor(0, 242, 254)
    step_shape.line.width = Pt(2)
    
    # Icon
    icon_box = slide3.shapes.add_textbox(Inches(x_pos + 0.1), Inches(y_pos + 0.1), Inches(0.8), Inches(0.8))
    icon_frame = icon_box.text_frame
    icon_para = icon_frame.paragraphs[0]
    icon_para.text = icon
    icon_para.font.size = Pt(40)
    
    # Title
    title_box_step = slide3.shapes.add_textbox(Inches(x_pos + 1.1), Inches(y_pos + 0.2), Inches(2.7), Inches(0.6))
    title_frame_step = title_box_step.text_frame
    title_para_step = title_frame_step.paragraphs[0]
    title_para_step.text = title
    title_para_step.font.size = Pt(22)
    title_para_step.font.bold = True
    title_para_step.font.color.rgb = RGBColor(79, 172, 254)
    
    # Description
    desc_box = slide3.shapes.add_textbox(Inches(x_pos + 0.15), Inches(y_pos + 1), Inches(3.7), Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_para = desc_frame.paragraphs[0]
    desc_para.text = desc
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = RGBColor(51, 51, 51)

# Principle box at bottom
principle_shape = slide3.shapes.add_shape(1, Inches(0.75), Inches(5.8), Inches(8.5), Inches(1.3))
principle_shape.fill.solid()
principle_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
principle_shape.line.color.rgb = RGBColor(51, 51, 51)
principle_shape.line.width = Pt(2)

principle_box = slide3.shapes.add_textbox(Inches(1), Inches(5.95), Inches(8), Inches(1))
principle_frame = principle_box.text_frame
principle_frame.word_wrap = True
principle_para = principle_frame.paragraphs[0]
principle_para.text = "💡 Key Principle: One concept = One file. Use [[brackets]] for links. Always cite sources. Document contradictions."
principle_para.font.size = Pt(18)
principle_para.font.bold = True
principle_para.font.color.rgb = RGBColor(51, 51, 51)
principle_para.alignment = PP_ALIGN.CENTER

# Speaker notes for slide 3
notes_slide3 = slide3.notes_slide
text_frame3 = notes_slide3.notes_text_frame
text_frame3.text = """The architecture is beautifully simple and follows a three-step cycle.

First, you give the system papers or raw materials. The AI compiles them—extracting key concepts and creating individual markdown files for each concept.

Second, when new sources arrive, the system syncs them. It updates existing pages with new information and creates new pages for novel concepts.

Third, the audit stage keeps quality high by finding orphan pages, missing references, and contradictions.

This creates a compounding feedback loop where knowledge gets richer and more interconnected over time.

The key principles ensure quality:
- One concept per file keeps things focused and reusable
- Bracket links create a knowledge graph
- Source citations maintain traceability
- Contradiction sections capture nuance"""

# ============= SLIDE 4: REAL-WORLD IMPLEMENTATION =============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.gradient()
fill.gradient_angle = 45.0
fill.gradient_stops[0].color.rgb = RGBColor(250, 112, 154)
fill.gradient_stops[1].color.rgb = RGBColor(254, 225, 64)

# Title
title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
title_frame4 = title_box4.text_frame
title_para4 = title_frame4.paragraphs[0]
title_para4.text = "🚀 Live Implementation"
title_para4.font.size = Pt(48)
title_para4.font.bold = True
title_para4.font.color.rgb = RGBColor(51, 51, 51)

# Stats boxes
stats = [
    ("3", "Copilot Skills"),
    ("100+", "Entity Pages"),
    ("⚙️", "Fully Automated"),
]

x_pos = 1.2
for stat_num, stat_label in stats:
    stat_box = slide4.shapes.add_textbox(Inches(x_pos), Inches(1.4), Inches(2.2), Inches(1.2))
    stat_frame = stat_box.text_frame
    stat_frame.word_wrap = True
    
    num_para = stat_frame.paragraphs[0]
    num_para.text = stat_num
    num_para.font.size = Pt(48)
    num_para.font.bold = True
    num_para.font.color.rgb = RGBColor(250, 112, 154)
    num_para.alignment = PP_ALIGN.CENTER
    
    label_para = stat_frame.add_paragraph()
    label_para.text = stat_label
    label_para.font.size = Pt(14)
    label_para.font.color.rgb = RGBColor(51, 51, 51)
    label_para.alignment = PP_ALIGN.CENTER
    label_para.space_before = Pt(6)
    
    x_pos += 2.8

# Features title
features_title = slide4.shapes.add_textbox(Inches(0.75), Inches(2.9), Inches(8.5), Inches(0.5))
features_title_frame = features_title.text_frame
features_title_para = features_title_frame.paragraphs[0]
features_title_para.text = "📦 What's Built"
features_title_para.font.size = Pt(28)
features_title_para.font.bold = True
features_title_para.font.color.rgb = RGBColor(51, 51, 51)

# Features list
features = [
    "✅ Automated concept extraction from PDFs & articles",
    "✅ Intelligent linking between related concepts",
    "✅ Incremental sync for new sources",
    "✅ Quality auditing at scale",
    "✅ Git-based version control",
]

y_pos = 3.6
for feature in features:
    feature_box = slide4.shapes.add_textbox(Inches(0.75), Inches(y_pos), Inches(8.5), Inches(0.5))
    feature_frame = feature_box.text_frame
    feature_para = feature_frame.paragraphs[0]
    feature_para.text = feature
    feature_para.font.size = Pt(18)
    feature_para.font.color.rgb = RGBColor(51, 51, 51)
    y_pos += 0.5

# Speaker notes for slide 4
notes_slide4 = slide4.notes_slide
text_frame4 = notes_slide4.notes_text_frame
text_frame4.text = """Let's see this in action. We've implemented the complete system using GitHub Copilot Skills.

You drop papers into the raw folder, run /Compile Papers to Wiki, and within minutes, you have structured entity pages with cross-links and source citations.

Add a new paper? Run /Sync Wiki from Raw and it intelligently updates existing pages while creating new ones.

Every 20 pages, you can audit quality to fix orphans and contradictions.

It's fully git-versioned, so you track every change—this is important for collaborative knowledge building.

The stats show:
- 3 custom Copilot commands/skills (Compile, Sync, Audit)
- Capable of generating 100+ entity pages from research materials
- Fully automated—minimal human intervention needed
- Saves hours of manual knowledge base management"""

# ============= SLIDE 5: DEMO & CALL TO ACTION =============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.gradient()
fill.gradient_angle = 45.0
fill.gradient_stops[0].color.rgb = RGBColor(48, 207, 208)
fill.gradient_stops[1].color.rgb = RGBColor(51, 8, 103)

# Title
title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
title_frame5 = title_box5.text_frame
title_para5 = title_frame5.paragraphs[0]
title_para5.text = "👉 See It In Action"
title_para5.font.size = Pt(48)
title_para5.font.bold = True
title_para5.font.color.rgb = RGBColor(255, 255, 255)

# Video placeholder
video_shape = slide5.shapes.add_shape(1, Inches(1.5), Inches(1.4), Inches(7), Inches(3.2))
video_shape.fill.solid()
video_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
video_shape.fill.transparency = 0.3
video_shape.line.color.rgb = RGBColor(255, 255, 255)
video_shape.line.width = Pt(3)

video_text_box = slide5.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(1.5))
video_text_frame = video_text_box.text_frame
video_text_frame.word_wrap = True

play_icon = video_text_frame.paragraphs[0]
play_icon.text = "▶️"
play_icon.font.size = Pt(60)
play_icon.alignment = PP_ALIGN.CENTER

video_label = video_text_frame.add_paragraph()
video_label.text = "Live Demo Video"
video_label.font.size = Pt(24)
video_label.font.bold = True
video_label.font.color.rgb = RGBColor(255, 255, 255)
video_label.alignment = PP_ALIGN.CENTER
video_label.space_before = Pt(10)

video_desc = video_text_frame.add_paragraph()
video_desc.text = "Watch the system compile a paper into entity pages in real-time"
video_desc.font.size = Pt(14)
video_desc.font.color.rgb = RGBColor(200, 200, 200)
video_desc.alignment = PP_ALIGN.CENTER
video_desc.space_before = Pt(6)

# CTA box
cta_shape = slide5.shapes.add_shape(1, Inches(0.75), Inches(4.8), Inches(8.5), Inches(2.2))
cta_shape.fill.solid()
cta_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
cta_shape.line.color.rgb = RGBColor(255, 255, 255)

cta_box = slide5.shapes.add_textbox(Inches(1), Inches(4.95), Inches(8), Inches(2))
cta_frame = cta_box.text_frame
cta_frame.word_wrap = True

cta_heading = cta_frame.paragraphs[0]
cta_heading.text = "🎯 Start Your Own Wiki Today"
cta_heading.font.size = Pt(24)
cta_heading.font.bold = True
cta_heading.font.color.rgb = RGBColor(48, 207, 208)

cta_repo = cta_frame.add_paragraph()
cta_repo.text = "📍 Repository: github.com/dev-enthusiast-84/llm-wiki"
cta_repo.font.size = Pt(14)
cta_repo.font.color.rgb = RGBColor(51, 51, 51)
cta_repo.space_before = Pt(8)

cta_start = cta_frame.add_paragraph()
cta_start.text = "🚀 Get Started: Clone → Add papers → Run /Compile Papers to Wiki"
cta_start.font.size = Pt(14)
cta_start.font.color.rgb = RGBColor(51, 51, 51)
cta_start.space_before = Pt(4)

cta_q = cta_frame.add_paragraph()
cta_q.text = "❓ Questions? Check the docs in skills/ folder"
cta_q.font.size = Pt(14)
cta_q.font.color.rgb = RGBColor(51, 51, 51)
cta_q.space_before = Pt(4)

# Speaker notes for slide 5
notes_slide5 = slide5.notes_slide
text_frame5 = notes_slide5.notes_text_frame
text_frame5.text = """This video shows the magic in action—a researcher adds research papers, triggers the compile command, and the system extracts concepts, creates pages, links related ideas, and generates an index.

In 10 minutes of computation, what might take a human a week is done. The best part? Every source is tracked, every concept is cited, and the knowledge base grows intelligently.

You can fork the repo, start with your own research, and build a compounding knowledge base. The wiki grows over time, and the more papers you add, the richer and more interconnected your knowledge becomes.

Call to action:
- Clone the repository
- Start with your research papers or the ones we've included
- Run the /Compile Papers to Wiki command
- Watch your wiki grow

This is how modern researchers should manage what they know. Thank you!"""

# Save presentation
prs.save('LLM-Wiki-Presentation.pptx')
print("✅ PowerPoint presentation created: LLM-Wiki-Presentation.pptx")
print("📊 5 slides with speaker notes, gradients, and icons")
print("⏱️  Optimized for 10-minute presentation")
