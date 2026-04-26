"""
LLM Wiki — Presentation Generator
Generates professional slides in PPTX, Keynote, and Google Slides formats.

Slide deck (5 slides, ~10 min):
  1. Title       — LLM Wiki: Karpathy's Knowledge Architecture
  2. Problem     — Why wikis fail; Karpathy's three-layer solution
  3. Architecture — Ingest · Query · Lint workflows + entity format
  4. Implementation — Claude Code CLI vs GitHub Copilot commands + Query UI
  5. Demo         — Embedded terminal and UI demo videos + CTA

Run:
    python generate_presentation.py             # all formats
    python generate_presentation.py --pptx      # PPTX only
    python generate_presentation.py --keynote   # Keynote only (macOS)
    python generate_presentation.py --slides    # Google Slides only

Requirements: pip install python-pptx
Optional:     pip install google-api-python-client google-auth-oauthlib
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE            = Path(__file__).parent
PPTX_OUT        = HERE / "llm-wiki-deck.pptx"
PPTX_LIGHT_OUT  = HERE / "llm-wiki-deck-light.pptx"
KEY_OUT         = HERE / "llm-wiki-deck.key"
TERMINAL_VID    = HERE / "terminal-demo.mp4"
UI_VID          = HERE / "ui-demo.mp4"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Color palette ──────────────────────────────────────────────────────────────

class C:
    BG        = (8,  12,  28)
    PANEL     = (16, 26,  56)
    CARD      = (22, 36,  72)
    DIVIDER   = (35, 55,  100)

    WHITE     = (240, 225, 195)   # warm oats cream — 14.4:1 on BG, 11.8:1 on CARD ✓
    LIGHT     = (196, 212, 238)
    DIM       = (105, 130, 172)

    BLUE      = ( 59, 130, 246)
    CYAN      = (  6, 182, 212)
    GREEN     = ( 16, 185, 129)
    ORANGE    = (245, 158,  11)
    PURPLE    = (139,  92, 246)
    RED       = (239,  68,  68)

    # AA-adjusted text colors for use on CARD background (22,36,72)
    # C.BLUE   on CARD = 4.14:1 (fail) → BLUE_AA = 6.06:1 (pass)
    # C.PURPLE on CARD = 3.59:1 (fail) → PURP_AA = 5.32:1 (pass)
    BLUE_AA   = (100, 165, 255)
    PURP_AA   = (170, 130, 255)

    # Gradient end colors (used by slide 1 and slide 5)
    GRAD1_END = (12,  22,  55)
    GRAD5_END = ( 5,  10,  32)


class CL:
    """Warm-oats light theme — WCAG AA verified on all warm backgrounds.

    Background luminances: BG≈0.941  PANEL≈0.872  CARD≈0.784
    All text contrast ratios against the background each colour touches:
      WHITE  on BG/PANEL/CARD  → 16.2 / 14.4 / 12.9 : 1  ✓
      LIGHT  on BG/PANEL/CARD  →  8.8 /  7.8 /  7.0 : 1  ✓
      DIM    on BG/PANEL/CARD  →  6.6 /  5.9 /  5.3 : 1  ✓ (normal text)
      BLUE   on BG/CARD        →  6.0 /  5.0 : 1          ✓
      CYAN   on BG/CARD        →  5.5 /  4.6 : 1          ✓
      GREEN  on BG/CARD (large)→  5.2 / 4.4:1 — only used ≥14pt bold (3:1)  ✓
      ORANGE on BG/CARD        →  5.7 /  4.8 : 1          ✓
      PURPLE on BG/CARD        →  7.0 /  5.9 : 1          ✓
      RED    on BG              →  6.2 : 1                 ✓
    """
    BG       = (252, 248, 240)   # warm oats background
    PANEL    = (244, 238, 226)   # warm oats panel
    CARD     = (236, 228, 212)   # warm oats card
    DIVIDER  = (180, 165, 140)   # warm divider

    WHITE    = ( 35,  26,  10)   # warm near-black primary text — 16.2:1 on BG ✓
    LIGHT    = ( 85,  68,  44)   # warm brown secondary text — 8.8:1 on BG, 7.0:1 on CARD ✓
    DIM      = (105,  86,  60)   # warm muted label text — 6.6:1 on BG, 5.3:1 on CARD ✓

    BLUE     = ( 20,  90, 200)   # 6.0:1 on BG, 5.0:1 on CARD ✓
    CYAN     = (  0, 110, 140)   # 5.5:1 on BG, 4.6:1 on CARD ✓
    GREEN    = (  0, 120,  75)   # 5.2:1 on BG; on CARD only at ≥14pt bold (3:1 needed) ✓
    ORANGE   = (160,  75,   0)   # 5.7:1 on BG, 4.8:1 on CARD ✓
    PURPLE   = (110,  40, 200)   # 7.0:1 on BG, 5.9:1 on CARD ✓
    RED      = (180,  30,  30)   # 6.2:1 on BG ✓

    # On warm CARD all dark accents already pass 4.5:1; same value as base colour
    BLUE_AA  = ( 20,  90, 200)
    PURP_AA  = (110,  40, 200)

    GRAD1_END = (248, 244, 234)   # warm gradient end for slide 1
    GRAD5_END = (250, 246, 236)   # warm gradient end for slide 5


# ── Primitive helpers ──────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg_grad(slide, c1, c2, angle=135):
    f = slide.background.fill
    f.gradient()
    f.gradient_angle = angle
    f.gradient_stops[0].color.rgb = RGBColor(*c1)
    f.gradient_stops[1].color.rgb = RGBColor(*c2)


def bg_solid(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = RGBColor(*color)


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, fill, border=None, bw=1.0):
    s = slide.shapes.add_shape(1, x, y, w, h)   # 1 = RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*fill)
    if border:
        s.line.color.rgb = RGBColor(*border)
        s.line.width = Pt(bw)
    else:
        _no_line(s)
    return s


def rrect(slide, x, y, w, h, fill, border=None, bw=1.0):
    s = slide.shapes.add_shape(5, x, y, w, h)   # 5 = ROUNDED_RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*fill)
    if border:
        s.line.color.rgb = RGBColor(*border)
        s.line.width = Pt(bw)
    else:
        _no_line(s)
    return s


def txt(slide, text: str, x, y, w, h, size=18, bold=False,
        color=None, align=PP_ALIGN.LEFT, font="Calibri", wrap=False):
    """color defaults to C.WHITE resolved at call time so theme-swap works correctly."""
    if color is None:
        color = C.WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size  = Pt(size)
    p.font.bold  = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name  = font
    p.alignment  = align


def txt_lines(slide, items, x, y, w, h, font="Calibri"):
    """items = [(text, size, bold, color, align), ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align) in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.font.size  = Pt(size)
        p.font.bold  = bold
        p.font.color.rgb = RGBColor(*color)
        p.font.name  = font
        p.alignment  = align


def pill(slide, label, x, y, w=Inches(1.8), text_color=C.CYAN, font_size=12):
    rrect(slide, x, y, w, Inches(0.36), C.PANEL, border=text_color, bw=0.8)
    txt(slide, label, x, y + Inches(0.04), w, Inches(0.28),
        size=font_size, color=text_color, align=PP_ALIGN.CENTER)


def hyperlink_pill(slide, label, url, x, y, w=Inches(1.8),
                   text_color=C.CYAN, font_size=14):
    """Styled pill button with a clickable, underlined hyperlink.
    14pt bold = WCAG 'large text' — 3:1 contrast required, all palette colors pass."""
    bg = tuple(max(0, min(255, int(c * 0.20))) for c in text_color)
    rrect(slide, x, y, w, Inches(0.46), bg, border=text_color, bw=1.5)
    tb = slide.shapes.add_textbox(x, y + Inches(0.08), w, Inches(0.30))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.underline = True        # AA: links distinguishable without colour alone
    run.font.color.rgb = RGBColor(*text_color)
    run.font.name = "Calibri"
    run.hyperlink.address = url


def vbar(slide, x, y, h, color=C.BLUE):
    rect(slide, x, y, Inches(0.055), h, color)


def hline(slide, x, y, w, color=C.DIVIDER):
    rect(slide, x, y, w, Pt(1), color)


def notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def video_or_placeholder(slide, video_path, x, y, w, h, label, color):
    """Embed video if file exists, otherwise render a styled placeholder."""
    if video_path and Path(video_path).exists():
        try:
            slide.shapes.add_movie(
                str(video_path), x, y, w, h, mime_type="video/mp4"
            )
            return
        except Exception:
            pass
    rrect(slide, x, y, w, h, C.CARD, border=color, bw=1.5)
    txt(slide, "▶", x, y + int(h * 0.28), w, Inches(0.9),
        size=52, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + int(h * 0.65), w, Inches(0.4),
        size=15, bold=True, color=C.LIGHT, align=PP_ALIGN.CENTER)
    txt(slide, "[ click to play embedded video ]", x, y + int(h * 0.8), w, Inches(0.32),
        size=11, color=C.LIGHT, align=PP_ALIGN.CENTER)  # DIM→LIGHT: 3.90→10.15:1 on CARD


# ── Slide 1: Title ─────────────────────────────────────────────────────────────

def slide_title(prs: Presentation):
    sl = _blank(prs)
    bg_grad(sl, C.BG, C.GRAD1_END, angle=135)

    # Left — hero text
    vbar(sl, Inches(0.55), Inches(1.85), Inches(3.2), color=C.BLUE)

    txt(sl, "LLM WIKI",
        Inches(0.75), Inches(1.75), Inches(8), Inches(1.35),
        size=76, bold=True, color=C.WHITE)

    txt(sl, "Karpathy's Knowledge Architecture for AI Research",
        Inches(0.75), Inches(3.2), Inches(7.6), Inches(0.75),
        size=24, color=C.LIGHT)

    # Attribute pills
    chips = [("AI-maintained", C.BLUE), ("Human-curated", C.CYAN), ("Always current", C.GREEN)]
    for i, (label, color) in enumerate(chips):
        pill(sl, label, Inches(0.75 + i * 2.0), Inches(4.15), w=Inches(1.85), text_color=color)

    txt(sl, "Inspired by Andrej Karpathy  ·  gist.github.com/karpathy/442a6bf...",
        Inches(0.75), Inches(4.72), Inches(8.5), Inches(0.4),
        size=13, color=C.DIM)

    # Right — Three-layer architecture diagram
    dx = Inches(9.35)
    rect(sl, dx - Inches(0.2), Inches(1.1), Inches(3.85), Inches(5.95), C.PANEL)
    txt(sl, "ARCHITECTURE", dx, Inches(1.22), Inches(3.45), Inches(0.38),
        size=11, bold=True, color=C.LIGHT, align=PP_ALIGN.CENTER)  # C.DIM→LIGHT: 4.39→11.41:1

    layers = [
        # (label, desc, accent_color, text_color_on_card)
        # accent_color used for border/bar (non-text, 3:1 sufficient)
        # text_color_on_card must pass 4.5:1 on CARD (22,36,72) for 13pt bold
        ("RAW SOURCES",     "Papers · HTML · PPTX · DOCX · XLSX",  C.ORANGE, C.ORANGE),
        ("WIKI LAYER",      "Entity pages · [[links]] · Citations", C.BLUE,   C.BLUE_AA),
        ("SCHEMA / CONFIG", "CLAUDE.md · Conventions · Workflows",  C.PURPLE, C.PURP_AA),
    ]
    ly = Inches(1.7)
    for i, (label, desc, accent, text_color) in enumerate(layers):
        rrect(sl, dx, ly, Inches(3.45), Inches(1.25), C.CARD, border=accent, bw=1.2)
        rect(sl, dx, ly, Inches(0.055), Inches(1.25), accent)
        txt(sl, label, dx + Inches(0.12), ly + Inches(0.1), Inches(3.2), Inches(0.42),
            size=13, bold=True, color=text_color)
        txt(sl, desc, dx + Inches(0.12), ly + Inches(0.58), Inches(3.2), Inches(0.5),
            size=12, color=C.LIGHT)
        if i < 2:
            txt(sl, "↓", dx + Inches(1.6), ly + Inches(1.25), Inches(0.26), Inches(0.38),
                size=15, color=C.LIGHT, align=PP_ALIGN.CENTER)  # DIM→LIGHT: 4.39→11.41:1 on PANEL
        ly += Inches(1.65)

    notes(sl, """SPEAKER NOTES — SLIDE 1: TITLE

Welcome everyone. Today I'm presenting LLM Wiki — a working implementation of Andrej Karpathy's vision for AI-powered knowledge management in research.

Key message: instead of manually maintaining a research wiki (which humans inevitably abandon), we let AI do the bookkeeping while humans stay in control of what goes in.

The system has three clean layers:
  • Raw Sources — immutable papers, HTML pages, spreadsheets that you add; the AI never modifies these
  • Wiki Layer — LLM-generated entity pages with cross-references; the AI owns and maintains this layer entirely
  • Schema / Config — the CLAUDE.md file that tells the AI how to structure everything

Over the next four slides we'll see why this matters, how it's built, and a live demo.

Timing: ~1.5 minutes on this slide.""")


# ── Slide 2: Problem & Insight ─────────────────────────────────────────────────

def slide_problem(prs: Presentation):
    sl = _blank(prs)
    bg_solid(sl, C.BG)

    txt(sl, "The Knowledge Management Crisis",
        Inches(0.55), Inches(0.45), Inches(12.3), Inches(0.8),
        size=34, bold=True, color=C.WHITE)
    hline(sl, Inches(0.55), Inches(1.28), Inches(12.25))

    cl = Inches(0.55)
    cr = Inches(7.0)
    cw = Inches(5.75)

    # Left — Problems
    txt(sl, "THE PROBLEM", cl, Inches(1.5), cw, Inches(0.38),
        size=12, bold=True, color=C.DIM)

    problems = [
        ("Knowledge Overload",
         "Hundreds of papers, articles, and docs — impossible to track manually"),
        ("Manual Maintenance",
         "Wikis die because nobody updates cross-references or fixes stale links"),
        ("Scattered Insights",
         "Key ideas live in separate notebooks, tabs, and memory — never connected"),
    ]
    py = Inches(2.0)
    for title, desc in problems:
        rrect(sl, cl, py, cw, Inches(1.32), C.CARD, border=C.RED, bw=0.7)
        txt(sl, title, cl + Inches(0.15), py + Inches(0.1), cw - Inches(0.2), Inches(0.42),
            size=15, bold=True, color=C.WHITE)
        txt(sl, desc, cl + Inches(0.15), py + Inches(0.56), cw - Inches(0.2), Inches(0.62),
            size=13, color=C.LIGHT, wrap=True)
        py += Inches(1.48)

    # Right — Karpathy's Insight
    txt(sl, "KARPATHY'S INSIGHT", cr, Inches(1.5), cw, Inches(0.38),
        size=12, bold=True, color=C.DIM)

    # Quote block
    rrect(sl, cr, Inches(2.0), cw, Inches(1.9), C.CARD, border=C.PURPLE, bw=1.5)
    txt(sl, '"LLMs don\'t get bored.\nThey don\'t forget to update a cross-reference."',
        cr + Inches(0.2), Inches(2.14), cw - Inches(0.35), Inches(1.38),
        size=16, bold=True, color=C.PURPLE)

    txt(sl, "Three-layer model:", cr, Inches(4.1), cw, Inches(0.4),
        size=14, bold=True, color=C.LIGHT)

    layer_items = [
        (C.ORANGE, "Raw Sources",  "Immutable — curated by humans"),
        (C.BLUE,   "Wiki Layer",   "LLM-owned — maintained automatically"),
        (C.PURPLE, "Schema",       "Config — conventions & workflows"),
    ]
    iy = Inches(4.6)
    for color, label, sub in layer_items:
        vbar(sl, cr, iy, Inches(0.55), color=color)
        txt(sl, label, cr + Inches(0.14), iy, Inches(2.1), Inches(0.3),
            size=14, bold=True, color=color)
        txt(sl, sub, cr + Inches(2.3), iy, Inches(3.3), Inches(0.3),
            size=13, color=C.LIGHT)
        iy += Inches(0.62)

    txt(sl, "Inspired by Vannevar Bush's Memex (1945) — the original vision for associative knowledge",
        cr, Inches(6.6), cw, Inches(0.45), size=11, color=C.DIM)

    notes(sl, """SPEAKER NOTES — SLIDE 2: THE PROBLEM

Every researcher and engineer faces this problem: we read papers, take notes, bookmark articles — but the knowledge stays scattered. Traditional wikis exist but nobody maintains them because the maintenance burden exceeds the perceived value.

Karpathy's key insight, inspired by Vannevar Bush's 1945 Memex paper, is simple: LLMs can be the maintenance layer. They don't tire of updating cross-references, don't forget to add citations, and can touch multiple files simultaneously.

The three-layer model separates concerns cleanly:
  • Raw Sources: humans curate what goes in — the AI never modifies these
  • Wiki Layer: the AI owns and maintains this — entity pages, links, summaries
  • Schema: the config file (CLAUDE.md) that tells the AI how to structure everything

The result: a wiki that actually stays current because the maintenance cost drops to near zero.

Timing: ~2 minutes on this slide.""")


# ── Slide 3: Architecture ──────────────────────────────────────────────────────

def slide_architecture(prs: Presentation):
    sl = _blank(prs)
    bg_solid(sl, C.BG)

    txt(sl, "Three Layers. Three Workflows.",
        Inches(0.55), Inches(0.45), Inches(12.3), Inches(0.8),
        size=34, bold=True, color=C.WHITE)
    hline(sl, Inches(0.55), Inches(1.28), Inches(12.25))

    # Workflow cards
    workflows = [
        ("INGEST",  C.BLUE,
         "Add a source → extract\nconcepts → update pages\n→ add cross-references"),
        ("QUERY",   C.CYAN,
         "Search relevant pages\n→ synthesize answer\n→ file new discoveries"),
        ("LINT",    C.GREEN,
         "Find orphan pages, gaps,\ncontradictions, and\nstale information"),
    ]

    box_w  = Inches(3.72)
    box_h  = Inches(2.85)
    box_y  = Inches(1.6)
    gap    = Inches(0.47)
    total  = 3 * box_w + 2 * gap
    start  = (SLIDE_W - total) / 2

    for i, (label, color, desc) in enumerate(workflows):
        bx = start + i * (box_w + gap)
        rrect(sl, bx, box_y, box_w, box_h, C.CARD, border=color, bw=1.5)
        rect(sl, bx, box_y, box_w, Inches(0.07), color)
        txt(sl, label, bx, box_y + Inches(0.14), box_w, Inches(0.52),
            size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(sl, desc, bx + Inches(0.2), box_y + Inches(0.82), box_w - Inches(0.4), Inches(1.85),
            size=14, color=C.LIGHT, wrap=True)
        if i < 2:
            ax = bx + box_w + Inches(0.06)
            txt(sl, "→", ax, box_y + Inches(1.1), Inches(0.36), Inches(0.55),
                size=22, color=C.DIM, align=PP_ALIGN.CENTER)

    # Entity page anatomy
    epy = Inches(4.65)
    txt(sl, "ENTITY PAGE FORMAT", Inches(0.55), epy, Inches(12), Inches(0.34),
        size=11, bold=True, color=C.DIM)

    sections = ["Definition", "Summary", "Explanation", "Related Concepts", "Sources", "Contradictions"]
    ax = Inches(0.55)
    for section in sections:
        pill(sl, section, ax, epy + Inches(0.4), w=Inches(1.95),
             text_color=C.CYAN, font_size=12)
        ax += Inches(2.1)

    # Supported formats
    fpy = Inches(5.75)
    txt(sl, "SOURCE FORMATS  (full text extraction)", Inches(0.55), fpy, Inches(12), Inches(0.34),
        size=11, bold=True, color=C.DIM)

    formats = [("PDF", C.GREEN), ("TXT", C.GREEN), ("HTML", C.GREEN),
               ("PPTX", C.CYAN), ("DOCX", C.CYAN), ("XLSX", C.CYAN)]
    fx = Inches(0.55)
    for fmt, color in formats:
        pill(sl, fmt + "  ✓", fx, fpy + Inches(0.4), w=Inches(1.95), text_color=color, font_size=12)
        fx += Inches(2.1)

    notes(sl, """SPEAKER NOTES — SLIDE 3: ARCHITECTURE

The architecture is elegant. Three workflows the AI runs:

INGEST: when you add a new source file, the AI reads it, extracts key concepts, creates or updates entity pages, and adds cross-references. Each concept gets its own markdown file.

QUERY: when you ask a question via the UI, the AI searches the relevant wiki pages, synthesizes an answer with citations, and can file valuable new insights back into the wiki.

LINT: periodically, the AI runs health checks — finding orphan pages (nothing links to them), missing pages (referenced with [[brackets]] but not created), contradictions between sources, and stale information.

Entity pages follow a consistent six-section format:
  Definition · Summary · Explanation · Related Concepts · Sources · Contradictions

The Sources section always cites paper, author, and page/section. Contradictions explicitly document when different sources disagree — this helps researchers understand how the field has evolved.

Six file formats are supported with full text extraction (PDF, TXT, HTML) or best-effort extraction (PPTX, DOCX, XLSX).

Timing: ~2 minutes on this slide.""")


# ── Slide 4: Implementation ────────────────────────────────────────────────────

def slide_implementation(prs: Presentation):
    sl = _blank(prs)
    bg_solid(sl, C.BG)

    txt(sl, "Works Where You Work",
        Inches(0.55), Inches(0.45), Inches(12.3), Inches(0.8),
        size=34, bold=True, color=C.WHITE)
    hline(sl, Inches(0.55), Inches(1.28), Inches(12.25))

    # AA-safe variants (defined in class C) for text on CARD background
    _BLUE_AA = C.BLUE_AA
    _PURP_AA = C.PURP_AA

    cl, cr, cw = Inches(0.55), Inches(7.0), Inches(5.72)
    col_top, col_h = Inches(1.5), Inches(4.45)   # col bottom = 5.95"
    # Layout constants — both columns use the same cmd_w/desc geometry.
    # cmd_w=2.65" fits "/regenerate-presentation" (24 Courier New 11pt mono chars, ~0.10"/char).
    # desc_w=2.65" fits the longest description (≤35 Calibri 11pt chars).
    row_sp   = Inches(0.33)   # row pitch: 7 wiki + divider + 3 maint fits in col_h=4.45"
    cmd_h    = Inches(0.28)   # single-line text box height for 11pt text
    cmd_w    = Inches(2.65)   # fits "/regenerate-presentation" comfortably @ Courier New 11pt
    desc_off = Inches(2.90)   # desc_x = col_x + desc_off  (= 0.15 left + 2.65 cmd + 0.10 gap)
    desc_w   = Inches(2.65)   # ≤ 35 Calibri 11pt chars (~0.06"/char avg) on one line

    # ── Claude Code CLI ────────────────────────────────────────────────────
    rrect(sl, cl, col_top, cw, col_h, C.CARD, border=C.ORANGE, bw=1.5)
    rect(sl, cl, col_top, cw, Inches(0.07), C.ORANGE)
    txt(sl, "CLAUDE CODE CLI", cl, col_top + Inches(0.07), cw, Inches(0.36),
        size=13, bold=True, color=C.ORANGE, align=PP_ALIGN.CENTER)

    txt(sl, "WIKI SKILLS", cl + Inches(0.15), col_top + Inches(0.48),
        Inches(2.0), Inches(0.22), size=11, bold=True, color=C.ORANGE)

    wiki_claude = [
        ("/orchestrate-wiki", "Full pipeline — parallel agents"),     # 31 chars
        ("/compile-papers",   "Extract concepts from /raw"),           # 26 chars
        ("/sync-wiki",        "Incremental update from /raw"),         # 28 chars
        ("/audit-wiki",       "Quality check — parallel agents"),      # 31 chars
        ("/reset-wiki",       "Wipe wiki for a fresh start"),          # 27 chars
        ("/launch-wiki-ui",   "Start Streamlit query interface"),      # 31 chars
        ("/stop-wiki-ui",     "Stop the running Streamlit app"),       # 30 chars
    ]
    cy = col_top + Inches(0.76)
    for cmd, desc in wiki_claude:
        txt(sl, cmd,  cl + Inches(0.15), cy, cmd_w,   cmd_h,
            size=11, bold=True, color=C.ORANGE, font="Courier New")
        txt(sl, desc, cl + desc_off,     cy, desc_w,  cmd_h,
            size=11, color=C.LIGHT)
        cy += row_sp

    cy += Inches(0.04)
    hline(sl, cl + Inches(0.15), cy, cw - Inches(0.3))
    cy += Inches(0.08)
    txt(sl, "REPO MAINTENANCE", cl + Inches(0.15), cy,
        Inches(2.2), Inches(0.22), size=11, bold=True, color=_PURP_AA)
    cy += Inches(0.22)

    maint_claude = [
        ("/sync-docs",               "Sync README, GitHub Pages, sub-docs"),  # 35 chars
        ("/run-maintenance",         "Tests, skill check, git health"),        # 30 chars
        ("/regenerate-presentation", "Rebuild demo videos and PPTX deck"),     # 33 chars
    ]
    for cmd, desc in maint_claude:
        txt(sl, cmd,  cl + Inches(0.15), cy, cmd_w,   cmd_h,
            size=11, bold=True, color=_PURP_AA, font="Courier New")
        txt(sl, desc, cl + desc_off,     cy, desc_w,  cmd_h,
            size=11, color=C.LIGHT)
        cy += row_sp

    # ── GitHub Copilot Chat ────────────────────────────────────────────────
    rrect(sl, cr, col_top, cw, col_h, C.CARD, border=C.BLUE, bw=1.5)
    rect(sl, cr, col_top, cw, Inches(0.07), C.BLUE)
    txt(sl, "GITHUB COPILOT CHAT", cr, col_top + Inches(0.07), cw, Inches(0.36),
        size=13, bold=True, color=_BLUE_AA, align=PP_ALIGN.CENTER)

    txt(sl, "WIKI SKILLS", cr + Inches(0.15), col_top + Inches(0.48),
        Inches(2.0), Inches(0.22), size=11, bold=True, color=_BLUE_AA)

    wiki_copilot = [
        ("Orchestrate Wiki",       "Full pipeline — parallel agents"),     # 31 chars
        ("Compile Papers to Wiki", "Extract from all /raw sources"),       # 29 chars
        ("Sync Wiki from Raw",     "Incremental update from /raw"),        # 28 chars
        ("Audit Wiki",             "Quality check — parallel agents"),     # 31 chars
        ("Reset Wiki",             "Wipe wiki for a fresh start"),         # 27 chars
        ("Launch Wiki UI",         "Start Streamlit query interface"),     # 31 chars
        ("Stop Wiki UI",           "Stop the running Streamlit app"),      # 30 chars
    ]
    cy = col_top + Inches(0.76)
    for cmd, desc in wiki_copilot:
        txt(sl, cmd,  cr + Inches(0.15), cy, cmd_w,   cmd_h,
            size=11, bold=True, color=_BLUE_AA, font="Courier New")
        txt(sl, desc, cr + desc_off,     cy, desc_w,  cmd_h,
            size=11, color=C.LIGHT)
        cy += row_sp

    cy += Inches(0.04)
    hline(sl, cr + Inches(0.15), cy, cw - Inches(0.3))
    cy += Inches(0.08)
    txt(sl, "REPO MAINTENANCE", cr + Inches(0.15), cy,
        Inches(2.2), Inches(0.22), size=11, bold=True, color=_PURP_AA)
    cy += Inches(0.22)

    maint_copilot = [
        ("Sync Docs",               "Sync README and documentation"),         # 29 chars
        ("Run Maintenance",         "Tests, skill check, git health"),         # 30 chars
        ("Regenerate Presentation", "Rebuild demo videos and PPTX deck"),      # 33 chars
    ]
    for cmd, desc in maint_copilot:
        txt(sl, cmd,  cr + Inches(0.15), cy, cmd_w,   cmd_h,
            size=11, bold=True, color=_PURP_AA, font="Courier New")
        txt(sl, desc, cr + desc_off,     cy, desc_w,  cmd_h,
            size=11, color=C.LIGHT)
        cy += row_sp

    # ── Streamlit Query UI band (y=6.05–7.13", bottom margin 0.37") ───────
    band_y, band_h = Inches(6.05), Inches(1.08)
    rrect(sl, cl, band_y, Inches(12.25), band_h, C.CARD, border=C.CYAN, bw=1.2)
    rect(sl, cl, band_y, Inches(0.07), band_h, C.CYAN)

    txt(sl, "STREAMLIT QUERY UI", cl + Inches(0.2), band_y + Inches(0.10),
        Inches(2.6), Inches(0.38), size=13, bold=True, color=C.CYAN)
    txt(sl, "Multi-provider · Grounded", cl + Inches(0.2), band_y + Inches(0.56),
        Inches(2.6), Inches(0.28), size=10, color=C.LIGHT)   # DIM→LIGHT: 10.14:1

    feat_items = [
        # sub-text capped at ~37 Calibri 11pt chars (~1.92") to fit 2.3" textbox on one line
        ("Multi-provider Q&A",        "Anthropic · OpenAI · Google · Mistral"),
        ("Wiki-grounded answers",      "No hallucination outside wiki content"),
        ("API key + model validation", "Session memory — never saved to disk"),
        ("Explicit 'not found'",       "Clear signal when topic is outside wiki"),
    ]
    fx = Inches(3.4)
    for title, sub in feat_items:
        txt(sl, "· " + title, fx, band_y + Inches(0.10), Inches(2.3), Inches(0.36),
            size=12, bold=True, color=C.LIGHT)
        txt(sl, sub, fx, band_y + Inches(0.54), Inches(2.3), Inches(0.36),
            size=11, color=C.LIGHT, wrap=False)               # wrap=False prevents overflow clipping
        fx += Inches(2.35)

    notes(sl, """SPEAKER NOTES — SLIDE 4: IMPLEMENTATION

The system works with both Claude Code CLI and GitHub Copilot Chat.
All skills now use parallel agents internally for significant speed gains.

WIKI SKILLS (orange = Claude / blue = Copilot) — the core wiki workflow:
• orchestrate-wiki — NEW: one command runs the full pipeline; detects wiki state, spawns parallel
  extraction agents (one per source file), runs 4 parallel audit agents, then syncs all docs
• compile-papers — reads all of /raw using parallel agents (one per file); creates wiki pages
• sync-wiki — detects new files; spawns parallel agents when multiple new sources exist
• audit-wiki — runs 4 parallel agents: orphan check, missing pages, contradictions, stale claims
• reset-wiki — wipe wiki for a fresh start with new topic area
• launch-wiki-ui — starts the Streamlit web app locally
• stop-wiki-ui — stops the running Streamlit process

REPO MAINTENANCE (purple, both tools) — separate from wiki operations:
• sync-docs — keeps README, GitHub Pages, and copilot-instructions in sync after code changes
• run-maintenance — spawns 4 parallel agents: tests, skill registration, wiki health, git status
• regenerate-presentation — rebuilds demo videos and this deck; detects changed scripts via git

GitHub Copilot Chat equivalents (same skill files, natural language — no slash prefix):
• Orchestrate Wiki, Compile Papers to Wiki, Sync Wiki from Raw, Audit Wiki, Reset Wiki, Launch Wiki UI, Stop Wiki UI
• Sync Docs, Run Maintenance, Regenerate Presentation

The Streamlit Query UI is a multi-provider web app:
  • Supports Anthropic, OpenAI, Google, Mistral, Cohere — enter any model ID + API key
  • Answers grounded strictly in wiki content (no hallucination outside wiki)
  • API key lives only in browser session memory — never written to disk
  • Explicit "No information found" when a topic isn't in the wiki — you always know your boundaries

Timing: ~2 minutes on this slide.""")


# ── Slide 5: Demo ──────────────────────────────────────────────────────────────

def slide_demo(prs: Presentation,
               terminal_video: Path | None = None,
               ui_video: Path | None = None):
    sl = _blank(prs)
    bg_grad(sl, C.BG, C.GRAD5_END, angle=135)

    txt(sl, "See It In Action",
        Inches(0.55), Inches(0.45), Inches(12.3), Inches(0.8),
        size=34, bold=True, color=C.WHITE)
    hline(sl, Inches(0.55), Inches(1.28), Inches(12.25))

    vid_w = Inches(5.85)
    vid_h = Inches(4.1)
    vid_y = Inches(1.5)

    # Terminal demo
    video_or_placeholder(sl, terminal_video, Inches(0.55), vid_y, vid_w, vid_h,
                         "Terminal Skills Demo", C.ORANGE)
    txt(sl, "Terminal Skills Demo",
        Inches(0.55), vid_y + vid_h + Inches(0.1), vid_w, Inches(0.36),
        size=13, bold=True, color=C.ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, "orchestrate-wiki (parallel agents)  ·  compile  ·  sync  ·  audit  ·  Copilot",
        Inches(0.55), vid_y + vid_h + Inches(0.46), vid_w, Inches(0.3),
        size=11, color=C.DIM, align=PP_ALIGN.CENTER)

    # UI demo
    video_or_placeholder(sl, ui_video, Inches(6.95), vid_y, vid_w, vid_h,
                         "Wiki Query UI Demo", C.CYAN)
    txt(sl, "Wiki Query UI Demo",
        Inches(6.95), vid_y + vid_h + Inches(0.1), vid_w, Inches(0.36),
        size=13, bold=True, color=C.CYAN, align=PP_ALIGN.CENTER)
    txt(sl, "settings  ·  happy-path query  ·  empty wiki  ·  not-found error",
        Inches(6.95), vid_y + vid_h + Inches(0.46), vid_w, Inches(0.3),
        size=11, color=C.DIM, align=PP_ALIGN.CENTER)

    # References footer — four centred hyperlink pills
    # widths: 2.18 + 1.58 + 1.58 + 1.90, gaps 0.22 each
    # total: 2.18 + 0.22 + 1.58 + 0.22 + 1.58 + 0.22 + 1.90 = 7.90"  →  start = (13.33 - 7.90) / 2 = 2.72"
    hyperlink_pill(sl, "↗  Karpathy's Gist",
                   "https://gist.github.com/karpathy/442a6bf555914893e9891c11519de94f",
                   Inches(2.72), Inches(6.55), w=Inches(2.18), text_color=C.PURPLE)
    hyperlink_pill(sl, "↗  Git Repo",
                   "https://github.com/dev-enthusiast-84/llm-wiki",
                   Inches(5.12), Inches(6.55), w=Inches(1.58), text_color=C.CYAN)
    hyperlink_pill(sl, "↗  Live Page",
                   "https://dev-enthusiast-84.github.io/llm-wiki/",
                   Inches(6.92), Inches(6.55), w=Inches(1.58), text_color=C.GREEN)
    hyperlink_pill(sl, "↗  LinkedIn Post",
                   "https://www.linkedin.com/posts/maneetta-antony_llm-wiki-share-7452168182221623296-acGl?utm_source=share&utm_medium=member_desktop&rcm=ACoAAAMKjaMBgVWqnGQ7U4godYqOh1rcqonVh74",
                   Inches(8.72), Inches(6.55), w=Inches(1.90), text_color=C.BLUE_AA)

    notes(sl, """SPEAKER NOTES — SLIDE 5: DEMO

[CLICK EACH VIDEO TO PLAY — both are embedded in the presentation]

TERMINAL DEMO covers:
  1. Claude Code CLI startup
  2. /compile-papers — scanning 9 /raw sources, creating 37 entity pages
  3. /sync-wiki — all 9 sources already compiled; wiki is up to date
  4. /audit-wiki — 37 pages, 252 links resolved, 0 orphans, 0 missing, 0 contradictions
  5. GitHub Copilot Chat equivalents (@workspace /Compile Papers to Wiki + /Launch Wiki UI)

UI DEMO covers:
  1. Settings screen — entering model ID and API key (happy path)
  2. Settings error — invalid API key format is caught immediately
  3. Empty wiki warning — clear instruction to run /compile-papers first
  4. Happy path query — "What is self-attention?" → cited answer from wiki
  5. Not-found error — "What is a diffusion model?" → explicit signal + how to add it

Key talking points:
  • Terminal: minutes from raw papers to 37 interconnected wiki pages
  • UI: the "not found" response is a feature, not a bug — you always know your wiki's boundaries
  • Both tools share the same wiki files — no duplication
  • Every change is git-versioned — safe to collaborate

After demo: "This is available today. Fork the repo, drop in your papers, run /compile-papers. Your wiki starts growing immediately."

Timing: ~2 minutes + demo playback.""")


# ── Build PPTX ─────────────────────────────────────────────────────────────────

def validate_layout(prs: Presentation) -> bool:
    """Check every shape stays within slide bounds. Returns True if clean."""
    sw, sh = prs.slide_width, prs.slide_height
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
            if None in (l, t, w, h):
                continue
            if l + w > sw:
                issues.append(f"Slide {i} '{shape.name}': right={(l+w)/914400:.3f}\" overflows slide_W={sw/914400:.3f}\"")
            if t + h > sh:
                issues.append(f"Slide {i} '{shape.name}': bottom={(t+h)/914400:.3f}\" overflows slide_H={sh/914400:.3f}\"")
            if l < 0:
                issues.append(f"Slide {i} '{shape.name}': left={l/914400:.3f}\" is off-canvas")
            if t < 0:
                issues.append(f"Slide {i} '{shape.name}': top={t/914400:.3f}\" is off-canvas")
    if issues:
        print(f"\n  ⚠  Layout issues ({len(issues)}):")
        for issue in issues:
            print(f"     {issue}")
    else:
        print("  ✓ Layout validation passed — all shapes within slide bounds")
    return len(issues) == 0


def build_pptx(
    output: Path | None = None,
    terminal_video: Path | None = TERMINAL_VID,
    ui_video: Path | None = UI_VID,
    theme=None,
) -> Path:
    """Generate the full 5-slide PPTX deck.

    theme: pass CL for the light theme; omit (or None) for the default dark theme.
    Swaps the module-level C reference so all slide functions pick up the right colours
    without any other changes.
    """
    global C
    old_C = C
    if theme is not None:
        C = theme
    if output is None:
        output = PPTX_LIGHT_OUT if theme is not None else PPTX_OUT
    try:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        slide_title(prs)
        slide_problem(prs)
        slide_architecture(prs)
        slide_implementation(prs)
        slide_demo(prs, terminal_video, ui_video)

        validate_layout(prs)
        prs.save(str(output))
        print(f"✓ PPTX  → {output}")
    finally:
        C = old_C
    return output


# ── Keynote conversion (macOS) ─────────────────────────────────────────────────

def convert_to_keynote(pptx_path: Path, key_path: Path = KEY_OUT) -> Path | None:
    """
    Open the PPTX in Keynote on macOS.

    Keynote 15+ on macOS Sequoia restricts programmatic PPTX→KEY export via
    AppleScript (sandbox prevents script-triggered file open). This function
    launches Keynote with the file and prints the one manual step needed.

    Returns key_path if it already exists from a prior manual export, else None.
    """
    if sys.platform != "darwin":
        print("⚠  Keynote conversion requires macOS")
        return None

    if key_path.exists():
        print(f"✓ Keynote → {key_path}  (already present)")
        return key_path

    # Detect installed Keynote app name
    app_name = None
    for name in ("Keynote Creator Studio", "Keynote"):
        probe = subprocess.run(
            ["osascript", "-e", f'tell application "{name}" to return version'],
            capture_output=True, text=True, timeout=5
        )
        if probe.returncode == 0 and probe.stdout.strip():
            app_name = name
            break

    if app_name is None:
        print("⚠  Keynote not found — install from Mac App Store")
        _keynote_manual_instructions(pptx_path, key_path)
        return None

    # Open in Keynote UI (PPTX opens natively; user confirms the conversion dialog)
    subprocess.run(["open", "-a", app_name, str(pptx_path.absolute())], check=False)
    print(f"  Launched {app_name} with the presentation.")
    _keynote_manual_instructions(pptx_path, key_path)
    return None


def _keynote_manual_instructions(pptx_path: Path, key_path: Path):
    print(f"\n  To save as Keynote (.key) format:")
    print(f"  1.  In Keynote → File → Export To → Keynote")
    print(f"  2.  Save to:  {key_path.absolute()}")
    print(f"  (The PPTX opens natively in Keynote — exporting to .key is optional.)")


# ── Google Slides ──────────────────────────────────────────────────────────────

def upload_to_google_slides(pptx_path: Path) -> str | None:
    """
    Upload PPTX to Google Drive, auto-converting to Google Slides.
    Returns the Google Slides URL on success, or None.

    Prerequisites:
        pip install google-api-python-client google-auth-oauthlib
        Place OAuth2 client-secret JSON at: ~/.llm-wiki-google-credentials.json
        (Google Cloud Console → APIs & Services → Credentials → OAuth 2.0 Client)
        Enable: Google Drive API + Google Slides API
    """
    try:
        from google.auth.transport.requests import Request
        from google.oauth2.credentials import Credentials
        from google_auth_oauthlib.flow import InstalledAppFlow
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload
    except ImportError:
        print("⚠  Google API client not installed.")
        print("   pip install google-api-python-client google-auth-oauthlib")
        _google_manual_instructions(pptx_path)
        return None

    SCOPES      = ["https://www.googleapis.com/auth/drive.file"]
    token_path  = Path.home() / ".llm-wiki-google-token.json"
    creds_path  = Path.home() / ".llm-wiki-google-credentials.json"

    creds = None
    if token_path.exists():
        creds = Credentials.from_authorized_user_file(str(token_path), SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        elif creds_path.exists():
            flow  = InstalledAppFlow.from_client_secrets_file(str(creds_path), SCOPES)
            creds = flow.run_local_server(port=0)
            token_path.write_text(creds.to_json())
        else:
            print(f"⚠  Credentials not found at {creds_path}")
            _google_manual_instructions(pptx_path)
            return None

    service = build("drive", "v3", credentials=creds)
    media   = MediaFileUpload(
        str(pptx_path),
        mimetype=(
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
    )
    metadata = {
        "name":     "LLM Wiki — Professional Presentation",
        "mimeType": "application/vnd.google-apps.presentation",
    }
    file = (
        service.files()
        .create(body=metadata, media_body=media, fields="id,webViewLink")
        .execute()
    )
    url = file.get("webViewLink", "")
    print(f"✓ Google Slides → {url}")
    return url


def _google_manual_instructions(pptx_path: Path):
    print("\nManual Google Slides import:")
    print(f"  1.  Open slides.google.com")
    print(f"  2.  File → Import slides → Upload")
    print(f"  3.  Select: {pptx_path.absolute()}")
    print("  Google auto-converts PPTX to Google Slides format.")


# ── CLI entry point ────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description="Generate LLM Wiki presentation")
    ap.add_argument("--pptx",    action="store_true", help="PPTX only")
    ap.add_argument("--keynote", action="store_true", help="Keynote (macOS)")
    ap.add_argument("--slides",  action="store_true", help="Google Slides upload")
    ap.add_argument("--light",   action="store_true", help="Light theme (WCAG AA)")
    args = ap.parse_args()

    all_fmt = not any([args.pptx, args.keynote, args.slides])
    theme   = CL if args.light else None

    pptx_path = None
    if all_fmt or args.pptx:
        pptx_path = build_pptx(theme=theme)

    default_pptx = PPTX_LIGHT_OUT if args.light else PPTX_OUT
    if all_fmt or args.keynote:
        if pptx_path is None:
            pptx_path = default_pptx
        convert_to_keynote(pptx_path)

    if all_fmt or args.slides:
        if pptx_path is None:
            pptx_path = default_pptx
        upload_to_google_slides(pptx_path)


if __name__ == "__main__":
    main()
